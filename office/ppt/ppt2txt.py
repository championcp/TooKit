from pptx import Presentation

def extract_text_from_ppt(ppt_file_path, output_text_file):
    prs = Presentation(ppt_file_path)

    with open(output_text_file, 'w', encoding='utf-8') as f:
        for slide_number, slide in enumerate(prs.slides):
            f.write(f"Slide {slide_number + 1}:\n")
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    f.write(shape.text + '\n')

ppt_file_path = 'your_presentation.pptx'
output_text_file = 'output_text.txt'
